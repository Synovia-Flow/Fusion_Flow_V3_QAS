"""Generate Fusion Flow V2 — Client MVP Onboarding presentation."""
import os, sys, subprocess

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

from pptx.oxml.ns import qn

# ── Colours (matching existing deck) ──────────────────────────────────────────
DARK_BLUE  = RGBColor(0x0A, 0x2A, 0x73)
MID_BLUE   = RGBColor(0x22, 0x67, 0xC6)
TEAL       = RGBColor(0x1A, 0xA6, 0xB8)
GREY       = RGBColor(0x6C, 0x7B, 0x97)
BODY       = RGBColor(0x2D, 0x3A, 0x55)
YELLOW     = RGBColor(0xF0, 0xB6, 0x4A)
RED        = RGBColor(0xDA, 0x4B, 0x5B)
GREEN      = RGBColor(0x27, 0xA3, 0x5D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF4, 0xF8, 0xFC)
TEAL_BG    = RGBColor(0xE8, 0xF6, 0xF8)
BLUE_BG    = RGBColor(0xCA, 0xD7, 0xEA)

W = Inches(13.333)
H = Inches(7.5)
SIDEBAR = Inches(0.22)

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "docs", "sop", "Fusion_Flow_V2_Client_MVP.pptx")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def add_rect(slide, l, t, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    solid_fill(shape, color)
    shape.line.fill.background()
    return shape

def add_label(slide, text, l, t, w, h, size, bold, color, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def section_tag(slide, text, color=TEAL):
    add_label(slide, text, Inches(0.45), Inches(0.32), Inches(5), Inches(0.35),
              12, True, color)

def slide_title(slide, title, subtitle=None):
    add_label(slide, title, Inches(0.45), Inches(0.70), Inches(12.5), Inches(0.9),
              26, True, DARK_BLUE)
    if subtitle:
        add_label(slide, subtitle, Inches(0.45), Inches(1.52), Inches(11), Inches(0.5),
                  16, False, GREY)

def sidebar(slide):
    add_rect(slide, 0, 0, SIDEBAR, H, DARK_BLUE)

def footer(slide, text="Synovia Digital / Fusion Flow V2"):
    add_label(slide, text, Inches(0.45), Inches(7.15), Inches(6), Inches(0.28),
              8, False, GREY)

def num_circle(slide, num, l, t, size=Inches(0.38)):
    circ = add_rect(slide, l, t, size, size, DARK_BLUE)
    circ.adjustments  # noqa
    add_label(slide, str(num), l, t, size, size, 12, True, WHITE, PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h, accent: RGBColor, title=None):
    bg_map = {
        TEAL:   RGBColor(0xE8,0xF6,0xF8),
        YELLOW: RGBColor(0xFE,0xF9,0xF1),
        RED:    RGBColor(0xFC,0xF2,0xF4),
        GREEN:  RGBColor(0xEF,0xF9,0xF4),
        MID_BLUE: RGBColor(0xEE,0xF3,0xFB),
    }
    bg = bg_map.get(accent, LIGHT_BG)
    add_rect(slide, l, t, w, h, bg).line.fill.background()
    # accent top bar
    add_rect(slide, l, t, w, Inches(0.06), accent).line.fill.background()
    cy = t + Inches(0.15)
    if title:
        add_label(slide, title, l + Inches(0.18), cy, w - Inches(0.25), Inches(0.35),
                  11, True, DARK_BLUE)
        cy += Inches(0.38)
    for item in items:
        # dot
        dot = slide.shapes.add_shape(9, l + Inches(0.18), cy + Inches(0.08),
                                     Inches(0.10), Inches(0.10))
        solid_fill(dot, accent)
        dot.line.fill.background()
        add_label(slide, item, l + Inches(0.36), cy, w - Inches(0.5), Inches(0.32),
                  10, False, BODY)
        cy += Inches(0.36)

def required_tag(slide, l, t):
    box = add_rect(slide, l, t, Inches(1.2), Inches(0.26), RED)
    box.line.fill.background()
    add_label(slide, "REQUIRED", l, t, Inches(1.2), Inches(0.26),
              8, True, WHITE, PP_ALIGN.CENTER)

def optional_tag(slide, l, t):
    box = add_rect(slide, l, t, Inches(1.2), Inches(0.26), GREEN)
    box.line.fill.background()
    add_label(slide, "OPTIONAL", l, t, Inches(1.2), Inches(0.26),
              8, True, WHITE, PP_ALIGN.CENTER)

def info_row(slide, label, description, req, l, t, w):
    # label box
    lw = Inches(2.8)
    lb = add_rect(slide, l, t, lw, Inches(0.42), LIGHT_BG)
    lb.line.fill.background()
    add_label(slide, label, l + Inches(0.1), t + Inches(0.08), lw - Inches(0.15),
              Inches(0.3), 9, True, DARK_BLUE)
    add_label(slide, description, l + lw + Inches(0.12), t + Inches(0.08),
              w - lw - Inches(1.45), Inches(0.3), 9, False, BODY)
    if req:
        required_tag(slide, l + w - Inches(1.25), t + Inches(0.08))
    else:
        optional_tag(slide, l + w - Inches(1.25), t + Inches(0.08))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
add_rect(sl, SIDEBAR, 0, W - SIDEBAR, H, BLUE_BG)

# Title block
add_rect(sl, Inches(0.45), Inches(1.6), Inches(8.5), Inches(0.06), TEAL).line.fill.background()
add_label(sl, "Fusion Flow V2 — Client MVP Setup", Inches(0.45), Inches(1.8),
          Inches(10), Inches(1.0), 30, True, DARK_BLUE)
add_label(sl, "What we need from you to get started",
          Inches(0.45), Inches(2.75), Inches(10), Inches(0.55), 20, False, GREY)

# Coloured chips
chips = [
    ("TSS API",    TEAL,     Inches(0.45)),
    ("Mailbox",    MID_BLUE, Inches(2.1)),
    ("Master Data",DARK_BLUE,Inches(3.75)),
    ("Transport",  YELLOW,   Inches(5.8)),
    ("SDI",        RED,      Inches(7.7)),
]
for label, col, x in chips:
    b = add_rect(sl, x, Inches(3.6), Inches(1.5), Inches(0.44), col)
    b.line.fill.background()
    tc = sl.shapes[-1]
    add_label(sl, label, x, Inches(3.6), Inches(1.5), Inches(0.44),
              10, True, WHITE, PP_ALIGN.CENTER)

add_label(sl, "Synovia Digital / Fusion Flow V2", Inches(0.45), Inches(7.0),
          Inches(6), Inches(0.3), 8, False, GREY)
add_label(sl, "Client onboarding requirements", Inches(0.45), Inches(7.22),
          Inches(6), Inches(0.3), 8, False, GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS THE MVP
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "OVERVIEW", TEAL)
slide_title(sl, "What You Get with the MVP",
            "A fully automated customs chain from inbound email to SupDec submission — with minimal manual intervention.")
footer(sl)

boxes = [
    (TEAL,     "📧  Inbound Email",    ["Graph mailbox polling", "Auto-classify attachments", "Excel sales orders parsed"]),
    (MID_BLUE, "🏗️  ENS Automation",   ["ENS header created", "Consignments + goods staged", "Pipeline validation runs"]),
    (DARK_BLUE,"🚀  TSS Submission",   ["Submit to TSS automatically", "DEC references stored back", "Status synced on schedule"]),
    (YELLOW,   "📄  SFD → SDI",        ["SFD detected after clearance", "SDI records auto-created", "SupDec submitted to TSS"]),
    (GREEN,    "🔔  Notifications",    ["Email on ENS received", "Movement authorised alert", "Staging failure alerts"]),
]
bw = Inches(2.35)
bh = Inches(3.5)
for i, (col, title, items) in enumerate(boxes):
    x = Inches(0.45) + i * (bw + Inches(0.18))
    bullet_box(sl, items, x, Inches(2.15), bw, bh, col, title)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT WE NEED: OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "REQUIREMENTS OVERVIEW", TEAL)
slide_title(sl, "Information We Need From You",
            "Five areas of data required before we can activate the full automation pipeline.")
footer(sl)

areas = [
    ("1", TEAL,     "TSS API Credentials",   "Your TSS API username and password to connect to the Trader Support Service."),
    ("2", MID_BLUE, "Microsoft Graph / Mailbox", "Azure App Registration credentials so we can poll your inbound email automatically."),
    ("3", DARK_BLUE,"Company & Master Data",  "Your EORI, company details, key partners and product catalogue."),
    ("4", YELLOW,   "Transport Defaults",     "Default carrier, arrival port, movement type and transport identity for your shipments."),
    ("5", RED,      "SDI / SupDec Defaults",  "Incoterms, representation type, postponed VAT and NI codes for SupDec automation."),
]

for i, (num, col, title, desc) in enumerate(areas):
    y = Inches(2.2) + i * Inches(0.88)
    num_circle(sl, num, Inches(0.45), y + Inches(0.02))
    add_rect(sl, Inches(0.92), y, Inches(11.8), Inches(0.7), LIGHT_BG).line.fill.background()
    add_rect(sl, Inches(0.92), y, Inches(0.06), Inches(0.7), col).line.fill.background()
    add_label(sl, title, Inches(1.1), y + Inches(0.08), Inches(3.2), Inches(0.3),
              11, True, DARK_BLUE)
    add_label(sl, desc, Inches(4.4), y + Inches(0.12), Inches(8.1), Inches(0.4),
              10, False, BODY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TSS API CREDENTIALS
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "AREA 1 OF 5 — TSS API", TEAL)
slide_title(sl, "TSS API Credentials",
            "Required to connect Fusion Flow V2 to the Trader Support Service REST API.")
footer(sl)

rows = [
    ("TSS API Username",   "Your TSS API account username (e.g. API.TSS0012045)",        True),
    ("TSS API Password",   "The password for the TSS API account",                        True),
    ("TSS Environment",    "Confirm: production or test (QAS). We default to production", True),
    ("ACT AS Customer ID", "customer_account_sys_id if using bureau/delegated access",    False),
]
y = Inches(2.2)
for label, desc, req in rows:
    info_row(sl, label, desc, req, Inches(0.45), y, Inches(12.5))
    y += Inches(0.52)

add_rect(sl, Inches(0.45), Inches(4.8), Inches(12.5), Inches(0.06), TEAL_BG).line.fill.background()
add_label(sl, "ℹ  Where to find this: TSS Portal → Account Settings → API Access. "
              "The API user must have the Declarant role assigned.",
          Inches(0.45), Inches(5.0), Inches(12.5), Inches(0.5), 9, False, GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MICROSOFT GRAPH / MAILBOX
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "AREA 2 OF 5 — MAILBOX", MID_BLUE)
slide_title(sl, "Microsoft Graph / Inbound Mailbox",
            "Allows Fusion Flow to poll your mailbox and pick up supplier emails automatically.")
footer(sl)

rows = [
    ("Azure Tenant ID",        "Your Microsoft 365 Azure AD tenant ID (GUID)",                   True),
    ("App Registration Client ID", "Client ID of the App Registration we will create together",  True),
    ("Client Secret",          "Secret generated from the App Registration",                     True),
    ("Mailbox Address",        "Email address of the shared mailbox to poll (e.g. ops@you.com)", True),
    ("Inbox Folder",           "Folder name to scan — default: INBOX. Subfolder: Inbox/BKD",    False),
    ("Processed Folder",       "Folder to move emails to after processing (e.g. Processed)",     False),
    ("Allowed Sender Domains", "Only process emails from these domains (e.g. supplier.com)",      False),
]
y = Inches(2.1)
for label, desc, req in rows:
    info_row(sl, label, desc, req, Inches(0.45), y, Inches(12.5))
    y += Inches(0.46)

add_label(sl, "ℹ  We can help you create the Azure App Registration. "
              "You will need an Azure AD admin to grant Mail.Read + Mail.ReadWrite permissions.",
          Inches(0.45), Inches(5.55), Inches(12.5), Inches(0.4), 9, False, GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — COMPANY & MASTER DATA
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "AREA 3 OF 5 — MASTER DATA", DARK_BLUE)
slide_title(sl, "Company & Master Data",
            "Reference data used to auto-populate and validate all declarations.")
footer(sl)

col1 = [
    ("Company EORI (GB)",   "Your GB EORI number (e.g. GB123456789000)",          True),
    ("Company EORI (XI)",   "Your XI EORI if trading into Northern Ireland",        False),
    ("Company Name",        "Trading name as registered with HMRC",                True),
    ("Company Address",     "Street, city, postcode, country",                     True),
]
col2 = [
    ("Key Partners",        "Main importers, exporters, carriers with EORI + address", True),
    ("Product Catalogue",   "Commodity codes (TARIC), gross weights, package types",   True),
    ("Haulier EORI",        "Default haulier EORI if different from carrier",          False),
    ("Consignor EORI",      "Default consignor EORI",                                  False),
]

y = Inches(2.15)
for label, desc, req in col1:
    info_row(sl, label, desc, req, Inches(0.45), y, Inches(6.1))
    y += Inches(0.52)

y = Inches(2.15)
for label, desc, req in col2:
    info_row(sl, label, desc, req, Inches(6.7), y, Inches(6.1))
    y += Inches(0.52)

add_label(sl, "ℹ  Partners and products can be bulk-imported via CSV. "
              "We provide templates.",
          Inches(0.45), Inches(4.55), Inches(12.5), Inches(0.35), 9, False, GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRANSPORT DEFAULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "AREA 4 OF 5 — TRANSPORT DEFAULTS", YELLOW)
slide_title(sl, "Transport & Carrier Defaults",
            "Applied automatically to every ENS when the email does not contain the value.")
footer(sl)

rows = [
    ("Arrival Port (UN/LOCODE)", "Port where goods arrive in the UK (e.g. GBHUL = Hull)",       True),
    ("Movement Type",            "Mode of transport: 1=Sea, 3=Road, 4=Air, 8=Inland waterway",   True),
    ("Transport Identity",       "Vessel IMO number or vehicle identity (e.g. IMO1234567)",       True),
    ("Transport Nationality",    "Country code of transport (e.g. GB)",                           True),
    ("Carrier EORI",             "EORI of the carrier company",                                   True),
    ("Carrier Name + Address",   "Name, street, city, postcode, country of the carrier",          True),
    ("Place of Loading",         "UN/LOCODE where goods are loaded",                              False),
    ("Place of Unloading",       "UN/LOCODE where goods are unloaded",                            False),
    ("Procedure Code",           "Default customs procedure code (e.g. 4000)",                    False),
    ("Invoice Currency",         "Default invoice currency code (e.g. GBP)",                     False),
]
y = Inches(2.1)
half = len(rows) // 2
for i, (label, desc, req) in enumerate(rows):
    col_x = Inches(0.45) if i < half else Inches(6.7)
    col_y = Inches(2.1) + (i % half) * Inches(0.48)
    info_row(sl, label, desc, req, col_x, col_y, Inches(6.1))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SDI / SUPDEC DEFAULTS
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "AREA 5 OF 5 — SDI DEFAULTS", RED)
slide_title(sl, "SDI / SupDec Automation Defaults",
            "Applied to every Supplementary Declaration created by the auto-discovery worker.")
footer(sl)

rows = [
    ("Incoterms",              "Default trade terms (e.g. DDP, CIF, FOB)",                      True),
    ("Representation Type",    "1=own account, 2=indirect, 3=direct (most common)",              True),
    ("Postponed VAT (PVA)",    "yes/no — whether to apply Postponed VAT Accounting by default",  True),
    ("Goods Domestic Status",  "D = domestic, N = non-domestic",                                 True),
    ("NI Additional Info Codes","NIREM for NI goods not at risk. Leave blank if not applicable", False),
    ("Nature of Transaction",  "Code for the nature of the commercial transaction (e.g. 11)",    False),
    ("Movement Type (SDI)",    "3 = RoRo/standard import",                                       False),
    ("Max Items per SupDec",   "Maximum goods items per SupDec before the system warns",         False),
]
y = Inches(2.1)
for label, desc, req in rows:
    info_row(sl, label, desc, req, Inches(0.45), y, Inches(12.5))
    y += Inches(0.50)

add_label(sl,
    "ℹ  The AUTOSUBMIT kill switch controls whether SupDecs are sent to TSS automatically. "
    "We keep this OFF until all defaults are confirmed and tested with you.",
    Inches(0.45), Inches(6.45), Inches(12.5), Inches(0.4), 9, False, GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "NEXT STEPS", GREEN)
slide_title(sl, "How We Get Started",
            "Once we receive the information, we can have the MVP running within a few days.")
footer(sl)

steps = [
    (GREEN,     "You provide",     [
        "TSS API credentials (username + password)",
        "Microsoft 365 App Registration (we guide you)",
        "Company EORI and address",
        "Partners and product list (CSV or Excel)",
        "Transport defaults for your typical shipments",
    ]),
    (TEAL,      "We configure",    [
        "Tenant provisioned in Fusion Flow V2",
        "TSS API connected and tested",
        "Graph mailbox polling activated",
        "Staging defaults set for your routes",
        "SDI automation configured (kill switch OFF initially)",
    ]),
    (MID_BLUE,  "We test together",[
        "Send a test email with a sample Sales Order",
        "Verify ENS is staged and validated",
        "Submit to TSS test environment",
        "Confirm SDI discovery works end-to-end",
        "Enable AUTOSUBMIT once all green",
    ]),
]
bw = Inches(3.9)
for i, (col, title, items) in enumerate(steps):
    x = Inches(0.45) + i * (bw + Inches(0.24))
    bullet_box(sl, items, x, Inches(2.15), bw, Inches(4.2), col, title)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CHECKLIST SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
sidebar(sl)
section_tag(sl, "SUMMARY CHECKLIST", TEAL)
slide_title(sl, "Client Information Checklist",
            "Use this as your reference when gathering information for the Fusion Flow V2 setup.")
footer(sl)

checklist = [
    (TEAL,     "TSS API",       ["TSS API username", "TSS API password", "Environment (production/test)", "ACT AS customer ID (if applicable)"]),
    (MID_BLUE, "Mailbox",       ["Azure Tenant ID", "App Client ID + Secret", "Mailbox email address", "Folder name + processed folder"]),
    (DARK_BLUE,"Master Data",   ["GB/XI EORI", "Company name + address", "Key partners with EORI", "Product catalogue (TARIC + weight)"]),
    (YELLOW,   "Transport",     ["Arrival port UN/LOCODE", "Movement type", "Carrier EORI + address", "Transport identity (IMO)"]),
    (RED,      "SDI",           ["Incoterms", "Representation type", "Postponed VAT yes/no", "NI additional info codes"]),
]

bw = Inches(2.35)
for i, (col, title, items) in enumerate(checklist):
    x = Inches(0.45) + i * (bw + Inches(0.18))
    bullet_box(sl, items, x, Inches(2.15), bw, Inches(3.8), col, title)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
